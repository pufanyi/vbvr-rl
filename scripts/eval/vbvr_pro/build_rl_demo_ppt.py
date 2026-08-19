#!/usr/bin/env python3
"""Build self-contained PowerPoint decks for the reviewed VBVR-Pro RL demos.

The generated presentations embed the native MP4 rollouts. Each movie poster is
the exact decoded first frame of that raw video, without labels or compositing.
"""

from __future__ import annotations

import argparse
import hashlib
import json
import zipfile
from collections import Counter
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from decord import VideoReader, cpu
from PIL import Image, ImageOps

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.oxml.ns import qn
    from pptx.oxml.xmlchemy import OxmlElement
    from pptx.util import Inches, Pt
except ModuleNotFoundError as exc:  # pragma: no cover - operator-facing dependency check
    raise SystemExit(
        "python-pptx is required. Install it in an isolated target and prepend that target to PYTHONPATH, e.g.\n"
        "  uv pip install --target storage/presentations/_tooling python-pptx==1.0.2"
    ) from exc


REPO_ROOT = Path(__file__).resolve().parents[3]
DEFAULT_OUTPUT_DIR = REPO_ROOT / "storage/presentations/vbvr_rl_demo_20260818"
SOURCE_PACKAGES = (
    (
        "A",
        "首批",
        REPO_ROOT / "storage/eval_out/vbvr_pro_rl_rollout_demo_final_strict_20260818",
    ),
    (
        "B",
        "新增",
        REPO_ROOT / "storage/eval_out/vbvr_pro_rl_rollout_demo_additional30_strict_20260818",
    ),
)

HERO_IDS = ("A01", "A02", "B01", "B02", "A07", "A09", "A06", "B06", "B09", "B12")

TASK_LABELS = {
    "G-131_select_next_figure_increasing_size_sequence_data-generator": "尺寸序列续项",
    "G-13_grid_number_sequence_data-generator": "数字序列路径",
    "G-15_grid_avoid_obstacles_data-generator": "避障网格",
    "G-161_mark_second_largest_shape_data-generator": "标记第二大图形",
    "G-16_grid_go_through_block_data-generator": "按序穿越色块",
    "G-202_mark_wave_peaks_data-generator": "标记波峰",
    "G-41_grid_highest_cost_data-generator": "最高成本路径",
    "G-45_key_door_matching_data-generator": "钥匙—门匹配",
    "G-47_multiple_keys_for_one_door_data-generator": "多钥匙单门",
    "G-54_connecting_color_data-generator": "同色连接",
    "O-16_color_addition_data-generator": "加色混合",
    "O-29_ballcolor_data-generator": "彩球合并与计数",
    "O-2_pigment_color_mixing_subtractive_data-generator": "减色混合",
    "O-31_ball_eating_data-generator": "球吞噬",
    "O-37_light_sequence_data-generator": "灯光序列",
    "O-39_maze_data-generator": "迷宫",
    "O-47_sliding_puzzle_data-generator": "滑块拼图",
    "O-5_symbol_deletion_data-generator": "符号删除",
    "O-75_communicating_vessels_data-generator": "连通器液面",
}

HERO_REASONS = {
    "A01": "四个离散选择覆盖四个候选，只有正确续项得满分。",
    "A02": "同一障碍网格出现四条不同路线，连续奖励从 0.120 到 0.997。",
    "B01": "第二大图形是需要全局比较的离散推理，四个选择只有一个正确。",
    "B02": "最终灯态与中间切换共同决定奖励，满分轨迹与 GT 一致。",
    "A07": "滑块移动序列与终盘都不同，奖励排序跟解题完成度一致。",
    "A09": "最高成本路径同时要求合法移动与全局优化，分数能区分路线质量。",
    "A06": "钥匙、门与终点形成多阶段规划，正确顺序的轨迹显著领先。",
    "B06": "彩球运动、合并和计数同时变化，规则奖励仍能给出合理排序。",
    "B09": "吞噬顺序影响动态过程和最终大小，四条行为轨迹本质不同。",
    "B12": "连通器需要持续物理演化，最终液面越接近平衡态得分越高。",
}

BG = "09111F"
PANEL = "111D30"
PANEL_2 = "17263B"
TEXT = "F4F7FB"
MUTED = "9FB0C5"
GRID = "2A3B52"
TEAL = "2DD4BF"
GREEN = "34D399"
BLUE = "60A5FA"
AMBER = "FBBF24"
ORANGE = "FB923C"
RED = "FB7185"
WHITE = "FFFFFF"
FONT_CN = "Noto Sans CJK SC"
FONT_LATIN = "Aptos"


@dataclass(frozen=True)
class DemoCase:
    package_code: str
    package_label: str
    root: Path
    data: dict[str, Any]

    @property
    def local_number(self) -> int:
        return int(self.data["display_id"].split("_")[-1])

    @property
    def case_id(self) -> str:
        return f"{self.package_code}{self.local_number:02d}"

    @property
    def task_label(self) -> str:
        return TASK_LABELS.get(self.data["task_name"], self.data["task_name"])

    @property
    def display_label(self) -> str:
        return f"{self.package_label} {self.local_number:02d}"

    def resolve(self, relative_path: str) -> Path:
        path = (self.root / relative_path).resolve()
        if not path.is_file():
            raise FileNotFoundError(path)
        return path


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--output-dir", type=Path, default=DEFAULT_OUTPUT_DIR)
    return parser.parse_args()


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


def set_run_typeface(run, typeface: str) -> None:
    """Set Latin, East Asian, and complex-script faces for portable CJK text."""
    run.font.name = typeface
    run_properties = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        element = run_properties.find(qn(tag))
        if element is None:
            element = OxmlElement(tag)
            run_properties.append(element)
        element.set("typeface", typeface)


def sha256_file(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for block in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest()


def load_cases() -> tuple[list[DemoCase], dict[str, Any]]:
    cases: list[DemoCase] = []
    manifests: dict[str, Any] = {}
    canonical_names: set[str] = set()
    for package_code, package_label, root in SOURCE_PACKAGES:
        manifest_path = root / "manifest.json"
        if not manifest_path.is_file():
            raise FileNotFoundError(manifest_path)
        manifest = json.loads(manifest_path.read_text())
        expected_count = 20 if package_code == "A" else 30
        if manifest["case_count"] != expected_count or len(manifest["cases"]) != expected_count:
            raise ValueError(f"Unexpected case count in {manifest_path}")
        manifests[package_code] = manifest
        for item in manifest["cases"]:
            canonical_name = item["canonical_name"]
            if canonical_name in canonical_names:
                raise ValueError(f"Duplicate canonical sample across packages: {canonical_name}")
            canonical_names.add(canonical_name)
            if item["manual_review"]["status"] != "pass":
                raise ValueError(f"Case is not manually approved: {package_code}{item['display_id']}")
            if len(item["rollouts"]) != 4 or item["diversity"]["unique_video_sha256"] != 4:
                raise ValueError(f"Case does not contain four unique rollouts: {package_code}{item['display_id']}")
            cases.append(DemoCase(package_code, package_label, root, item))
    if len(cases) != 50:
        raise ValueError(f"Expected 50 cases, found {len(cases)}")
    return cases, manifests


def make_raw_video_poster(video_path: Path, output_path: Path) -> None:
    """Write the raw video's decoded frame zero losslessly, without modification."""
    reader = VideoReader(str(video_path), ctx=cpu(0), num_threads=1)
    if len(reader) != 81:
        raise ValueError(f"Expected 81 frames in {video_path}, found {len(reader)}")
    image = Image.fromarray(reader[0].asnumpy()).convert("RGB")
    output_path.parent.mkdir(parents=True, exist_ok=True)
    image.save(output_path, format="PNG", optimize=True)


def extract_final_frame(video_path: Path, output_path: Path, size: int = 512) -> None:
    if output_path.is_file():
        return
    reader = VideoReader(str(video_path), ctx=cpu(0), num_threads=1)
    image = Image.fromarray(reader[-1].asnumpy()).convert("RGB")
    image = ImageOps.fit(image, (size, size), method=Image.Resampling.LANCZOS)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    image.save(output_path, format="PNG", optimize=True)


def prepare_posters(cases: list[DemoCase], output_dir: Path) -> dict[str, dict[str, Path]]:
    poster_root = output_dir / "assets/posters"
    assets: dict[str, dict[str, Path]] = {}
    for case in cases:
        case_assets: dict[str, Path] = {}
        for rollout in case.data["rollouts"]:
            index = rollout["rollout_index"]
            source = case.resolve(rollout["native_video"])
            poster = poster_root / case.case_id / f"rollout_{index:02d}.png"
            make_raw_video_poster(source, poster)
            case_assets[f"rollout_{index:02d}"] = poster
        gt_poster = poster_root / case.case_id / "ground_truth_final.png"
        extract_final_frame(case.resolve(case.data["ground_truth_video"]), gt_poster)
        case_assets["ground_truth_final"] = gt_poster
        assets[case.case_id] = case_assets
    return assets


def new_presentation() -> Presentation:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333333)
    presentation.slide_height = Inches(7.5)
    presentation.core_properties.title = "VBVR-Pro RL four-rollout demo"
    presentation.core_properties.subject = "DanceGRPO checkpoint evidence with embedded MP4 rollouts"
    presentation.core_properties.author = "VBVR-RL"
    presentation.core_properties.keywords = "RL, DanceGRPO, VBVR-Pro, CPS, embedded video"
    return presentation


def blank_slide(presentation: Presentation):
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(BG)
    return slide


def add_rect(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    fill_color: str,
    *,
    line_color: str | None = None,
    radius: bool = True,
    line_width: float = 1.0,
):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_color)
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = rgb(line_color)
        shape.line.width = Pt(line_width)
    return shape


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 18,
    color: str = TEXT,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.MIDDLE,
    font: str = FONT_CN,
    margin: float = 0.02,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(margin)
    frame.margin_right = Inches(margin)
    frame.margin_top = Inches(margin)
    frame.margin_bottom = Inches(margin)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    paragraph.space_before = Pt(0)
    paragraph.space_after = Pt(0)
    run = paragraph.add_run()
    run.text = text
    set_run_typeface(run, font)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return box


def add_multiline(
    slide,
    lines: list[tuple[str, float, str, bool]],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    bullet: bool = False,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.03)
    frame.margin_right = Inches(0.03)
    frame.margin_top = Inches(0.03)
    frame.margin_bottom = Inches(0.03)
    for index, (line, size, color, bold) in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.space_before = Pt(0)
        paragraph.space_after = Pt(5)
        paragraph.level = 0
        if bullet:
            paragraph.text = f"•  {line}"
            run = paragraph.runs[0]
        else:
            run = paragraph.add_run()
            run.text = line
        set_run_typeface(run, FONT_CN)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
    return box


def add_slide_chrome(slide, title: str, subtitle: str, slide_number: int) -> None:
    add_rect(slide, 0.0, 0.0, 0.08, 7.5, TEAL, radius=False)
    add_text(slide, title, 0.42, 0.25, 11.9, 0.48, size=26, bold=True)
    add_text(slide, subtitle, 0.44, 0.77, 11.8, 0.25, size=10.5, color=MUTED, font=FONT_LATIN)
    add_text(slide, f"{slide_number:02d}", 12.5, 0.3, 0.42, 0.28, size=10, color=MUTED, align=PP_ALIGN.RIGHT)


def add_notes(slide, text: str) -> None:
    notes_frame = slide.notes_slide.notes_text_frame
    notes_frame.text = text
    for paragraph in notes_frame.paragraphs:
        for run in paragraph.runs:
            set_run_typeface(run, FONT_CN)
            run.font.size = Pt(11)


def score_color(score: float) -> str:
    if score >= 0.85:
        return GREEN
    if score >= 0.6:
        return TEAL
    if score >= 0.3:
        return AMBER
    return RED


def competition_ranks(scores: list[float]) -> list[int]:
    ordered = sorted(scores, reverse=True)
    return [ordered.index(score) + 1 for score in scores]


def add_cover_slide(presentation: Presentation, hero_cases: list[DemoCase], assets: dict[str, dict[str, Path]]) -> None:
    slide = blank_slide(presentation)
    add_rect(slide, 0.0, 0.0, 0.11, 7.5, TEAL, radius=False)
    add_text(slide, "RL 真正在学什么？", 0.65, 0.65, 7.3, 0.7, size=34, bold=True)
    add_text(slide, "同一输入 · 四种行为 · 可验证的奖励差", 0.67, 1.42, 10.8, 0.48, size=24, color=TEAL, bold=True)
    add_text(
        slide,
        "Wan2.2-TI2V-5B  ·  DanceGRPO  ·  VBVR-Pro  ·  Flow-CPS 0.7",
        0.68,
        2.05,
        10.5,
        0.32,
        size=13,
        color=MUTED,
        font=FONT_LATIN,
    )
    add_rect(slide, 0.67, 2.62, 4.3, 0.04, TEAL, radius=False)
    add_text(slide, "50 组任务", 0.7, 2.92, 2.2, 0.45, size=25, bold=True)
    add_text(slide, "200 个内嵌 MP4 rollout", 2.9, 2.92, 4.2, 0.45, size=25, bold=True)
    add_text(slide, "8 个中间 checkpoint · 19 类任务", 0.7, 3.42, 6.7, 0.3, size=15, color=MUTED)
    for index, case in enumerate(hero_cases[:4]):
        x = 0.7 + index * 3.1
        add_rect(slide, x - 0.04, 4.18, 2.82, 2.52, PANEL, line_color=GRID, radius=True)
        slide.shapes.add_picture(
            str(assets[case.case_id]["rollout_00"]),
            Inches(x),
            Inches(4.22),
            Inches(2.74),
            Inches(2.0),
        )
        add_text(
            slide,
            f"{case.task_label}  ·  ckpt {case.data['checkpoint']}",
            x,
            6.27,
            2.74,
            0.28,
            size=10.5,
            bold=True,
            align=PP_ALIGN.CENTER,
        )
    add_text(slide, "2026-08-18", 10.9, 7.05, 1.8, 0.2, size=9.5, color=MUTED, align=PP_ALIGN.RIGHT)
    add_notes(
        slide,
        "开场：这套材料分成两层证据。完整 500 样本 checkpoint 曲线回答‘RL 整体有没有提升’，"
        "同一输入的四路 rollout 回答‘GRPO 在一个训练组里到底看到了什么学习信号’。",
    )


def add_summary_slide(presentation: Presentation, manifests: dict[str, Any]) -> None:
    slide = blank_slide(presentation)
    add_slide_chrome(slide, "结论先行", "QUANTITATIVE + QUALITATIVE EVIDENCE", len(presentation.slides))
    curve = manifests["A"]["aggregate_evidence"]["cps_0p7"]
    relative_gain = curve["best_delta"] / curve["baseline_overall"]
    metrics = (
        (f"+{curve['best_delta']:.4f}", "最佳 checkpoint 的绝对增益", TEAL),
        (f"+{relative_gain:.1%}", "相对 matched baseline", GREEN),
        ("50 / 200", "任务组 / rollout 视频", BLUE),
        ("200 / 200", "独立复算 exact match", AMBER),
    )
    for index, (value, label, color) in enumerate(metrics):
        x = 0.48 + index * 3.16
        add_rect(slide, x, 1.35, 2.88, 1.35, PANEL, line_color=GRID)
        add_text(slide, value, x + 0.16, 1.53, 2.56, 0.5, size=29, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, label, x + 0.16, 2.12, 2.56, 0.28, size=11.5, color=MUTED, align=PP_ALIGN.CENTER)
    add_rect(slide, 0.48, 3.08, 12.35, 3.7, PANEL, line_color=GRID)
    lines = [
        ("整体结果：CPS 0.7 从 0.472177 提升到 0.547886，最佳点在 checkpoint 2200。", 18, TEXT, True),
        ("组内行为：同一 prompt / input / checkpoint 下，四个随机种子产生语义上不同的决策或轨迹。", 17, TEXT, False),
        ("奖励可信：每个案例都人工对照 GT；最终两批分别 80/80、120/120 独立重算完全一致。", 17, TEXT, False),
        ("展示边界：50 组是为可解释性挑选的高对比案例；它们说明学习信号，不替代完整曲线。", 17, AMBER, False),
    ]
    add_multiline(slide, lines, 0.85, 3.48, 11.6, 2.85, bullet=True)
    add_notes(
        slide,
        "先给定量结论，再说明案例的角色。不要用挑选出的 50 组估计平均模型质量；"
        "平均提升由每格 500 个样本的固定合同曲线支撑。",
    )


def add_pipeline_slide(presentation: Presentation) -> None:
    slide = blank_slide(presentation)
    add_slide_chrome(
        slide, "GRPO 的组内学习信号从哪里来", "SAME CONDITION, FOUR STOCHASTIC OUTCOMES", len(presentation.slides)
    )
    boxes = (
        ("同一条件", "prompt + input\ncheckpoint k", BLUE),
        ("四次 rollout", "G = 4\nCPS = 0.7", TEAL),
        ("规则奖励", "r₁, r₂, r₃, r₄\nEvalKit", AMBER),
        ("组内优势", "高分 ↑\n低分 ↓", ORANGE),
        ("更新策略", "提高优质行为\n的相对概率", GREEN),
    )
    for index, (title, body, color) in enumerate(boxes):
        x = 0.45 + index * 2.55
        add_rect(slide, x, 2.0, 2.15, 2.15, PANEL, line_color=color, line_width=2.0)
        add_rect(slide, x, 2.0, 2.15, 0.12, color, radius=False)
        add_text(slide, title, x + 0.12, 2.35, 1.91, 0.38, size=18, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, body, x + 0.12, 2.86, 1.91, 0.72, size=15, align=PP_ALIGN.CENTER)
        if index < len(boxes) - 1:
            add_text(slide, "→", x + 2.18, 2.78, 0.34, 0.4, size=26, color=MUTED, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 0.72, 4.75, 11.85, 1.35, PANEL_2, line_color=GRID)
    add_text(
        slide,
        "案例页直接展示前三步：条件相同，但行为本质不同，并且奖励排序能被视觉结果解释。",
        1.0,
        4.98,
        11.3,
        0.4,
        size=19,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        "完整 checkpoint 曲线验证第五步确实转化成了总体得分提升。",
        1.0,
        5.47,
        11.3,
        0.3,
        size=15,
        color=TEAL,
        align=PP_ALIGN.CENTER,
    )
    add_notes(slide, "这里解释案例与训练的联系：案例不是 cherry-pick 后的平均性能，而是组内 advantage 的可视化。")


def add_curve_slide(presentation: Presentation, manifests: dict[str, Any]) -> None:
    slide = blank_slide(presentation)
    add_slide_chrome(
        slide,
        "定量证据：RL checkpoint 曲线整体上升",
        "500 SAMPLES PER CELL · 144 COMPLETE CELLS · 0 SCORER ERRORS",
        len(presentation.slides),
    )
    plot_path = SOURCE_PACKAGES[0][2] / manifests["A"]["aggregate_evidence"]["plot_png"]
    add_rect(slide, 0.42, 1.18, 12.45, 4.62, WHITE, line_color=GRID)
    slide.shapes.add_picture(str(plot_path), Inches(0.52), Inches(1.31), Inches(12.25), Inches(4.36))
    add_rect(slide, 0.7, 6.08, 3.75, 0.72, PANEL_2, line_color=TEAL)
    add_text(slide, "0.4722 → 0.5479", 0.9, 6.2, 3.35, 0.28, size=19, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_rect(slide, 4.78, 6.08, 3.75, 0.72, PANEL_2, line_color=GREEN)
    add_text(
        slide, "best @ checkpoint 2200", 4.98, 6.2, 3.35, 0.28, size=17, color=GREEN, bold=True, align=PP_ALIGN.CENTER
    )
    add_rect(slide, 8.87, 6.08, 3.75, 0.72, PANEL_2, line_color=AMBER)
    add_text(
        slide, "latest 2300: +0.0712", 9.07, 6.2, 3.35, 0.28, size=17, color=AMBER, bold=True, align=PP_ALIGN.CENTER
    )
    add_notes(
        slide,
        "定量主张：matched DiffSynth step-35500 baseline 在 CPS 0.7 下为 0.472177；"
        "checkpoint 2200 为 0.547886，绝对提升 0.075709，约 16.0%。最新 2300 仍高 0.071235。",
    )


def add_reading_guide_slide(
    presentation: Presentation,
    example: DemoCase,
    assets: dict[str, dict[str, Path]],
) -> None:
    slide = blank_slide(presentation)
    add_slide_chrome(slide, "案例页怎么读", "CLICK A POSTER TO PLAY THE EMBEDDED MP4", len(presentation.slides))
    scores = example.data["scores"]
    ranks = competition_ranks(scores)
    for index, rollout in enumerate(example.data["rollouts"]):
        x = 0.47 + index * 3.18
        add_text(
            slide,
            f"rollout {index} · seed {rollout['seed']}",
            x,
            1.25,
            2.82,
            0.25,
            size=10.5,
            color=MUTED,
            align=PP_ALIGN.CENTER,
        )
        add_rect(slide, x - 0.04, 1.54, 2.9, 2.9, PANEL, line_color=score_color(scores[index]), line_width=2.0)
        slide.shapes.add_picture(
            str(assets[example.case_id][f"rollout_{index:02d}"]), Inches(x), Inches(1.58), Inches(2.82), Inches(2.82)
        )
        add_text(
            slide,
            f"reward {scores[index]:.3f}   ·   #{ranks[index]}",
            x,
            4.55,
            2.82,
            0.36,
            size=17,
            color=score_color(scores[index]),
            bold=True,
            align=PP_ALIGN.CENTER,
        )
    guide = [
        ("每个大画面实际是内嵌 MP4；单击即可从 t=0 播放。", 15.5, TEXT, True),
        ("播放前直接显示 raw video 的原始第 0 帧；不加标注、不拼图。", 15, TEXT, False),
        ("四路只改变 seed；prompt、输入、checkpoint、CPS=0.7 和推理合同完全相同。", 15, TEXT, False),
        ("颜色表示绝对奖励区间，#1–#4 表示该 GRPO 组内排序。", 15, TEAL, False),
    ]
    add_rect(slide, 0.72, 5.22, 11.9, 1.38, PANEL_2, line_color=GRID)
    add_multiline(slide, guide, 0.95, 5.39, 11.4, 1.05, bullet=True)
    add_notes(slide, "演示时点击任意一个 raw 首帧封面即可播放对应 MP4。建议先播最低分，再播最高分，方便解释奖励排序。")


def add_case_slide(
    presentation: Presentation,
    case: DemoCase,
    assets: dict[str, dict[str, Path]],
    *,
    hero_index: int | None = None,
) -> None:
    slide = blank_slide(presentation)
    title_prefix = f"主讲 {hero_index:02d}/10" if hero_index is not None else case.display_label
    title = f"{title_prefix} · checkpoint {case.data['checkpoint']} · {case.task_label}"
    domain = "In-Domain" if case.data["domain"] == "In_Domain" else "Out-of-Domain"
    subtitle = (
        f"{case.case_id}  |  {domain}  |  CPS 0.7  |  score range {case.data['score_range']:.3f}  |  "
        f"{case.data['task_name']}"
    )
    add_slide_chrome(slide, title, subtitle, len(presentation.slides))
    scores = [float(score) for score in case.data["scores"]]
    ranks = competition_ranks(scores)
    for index, rollout in enumerate(case.data["rollouts"]):
        x = 0.44 + index * 3.18
        add_text(
            slide,
            f"rollout {index}  ·  seed {rollout['seed']}",
            x,
            1.12,
            2.84,
            0.23,
            size=10.2,
            color=MUTED,
            align=PP_ALIGN.CENTER,
            font=FONT_LATIN,
        )
        add_rect(slide, x - 0.035, 1.39, 2.91, 2.91, PANEL, line_color=score_color(scores[index]), line_width=2.0)
        movie_path = case.resolve(rollout["native_video"])
        poster_path = assets[case.case_id][f"rollout_{index:02d}"]
        slide.shapes.add_movie(
            str(movie_path),
            Inches(x),
            Inches(1.425),
            Inches(2.84),
            Inches(2.84),
            str(poster_path),
            mime_type="video/mp4",
        )
        add_text(
            slide,
            f"{scores[index]:.3f}",
            x + 0.18,
            4.43,
            1.35,
            0.38,
            size=21,
            color=score_color(scores[index]),
            bold=True,
            align=PP_ALIGN.LEFT,
            font=FONT_LATIN,
        )
        add_text(
            slide,
            f"组内 #{ranks[index]}",
            x + 1.48,
            4.45,
            1.18,
            0.33,
            size=11.5,
            color=MUTED,
            bold=True,
            align=PP_ALIGN.RIGHT,
        )
    add_text(slide, "输入", 0.44, 5.06, 1.22, 0.22, size=10.5, color=MUTED, align=PP_ALIGN.CENTER)
    slide.shapes.add_picture(
        str(case.resolve(case.data["input_image"])), Inches(0.44), Inches(5.31), Inches(1.22), Inches(1.22)
    )
    add_text(slide, "GT 最终帧", 1.8, 5.06, 1.22, 0.22, size=10.5, color=MUTED, align=PP_ALIGN.CENTER)
    slide.shapes.add_picture(
        str(assets[case.case_id]["ground_truth_final"]),
        Inches(1.8),
        Inches(5.31),
        Inches(1.22),
        Inches(1.22),
    )
    add_rect(slide, 3.25, 5.07, 9.62, 1.48, PANEL_2, line_color=GRID)
    add_text(slide, "人工复核  ✓ PASS", 3.48, 5.24, 2.35, 0.25, size=12, color=TEAL, bold=True)
    note = case.data["manual_review"]["note"]
    add_text(slide, note, 3.48, 5.55, 9.04, 0.72, size=13.2, valign=MSO_ANCHOR.TOP)
    if hero_index is not None:
        add_text(slide, HERO_REASONS[case.case_id], 3.48, 6.28, 8.85, 0.22, size=10.2, color=AMBER)
    package_exact = "80/80" if case.package_code == "A" else "120/120"
    add_text(
        slide,
        f"点击视频播放  ·  native 512×512×81 @ 16 FPS  ·  该批独立复算 {package_exact} exact match",
        0.48,
        6.86,
        12.2,
        0.2,
        size=9.5,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )
    rollout_lines = "\n".join(
        f"rollout {rollout['rollout_index']}: seed={rollout['seed']}, reward={rollout['score']:.9f}, "
        f"sha256={rollout['native_sha256']}"
        for rollout in case.data["rollouts"]
    )
    add_notes(
        slide,
        f"{case.display_label} / {case.case_id}\n"
        f"Task: {case.data['task_name']}\n"
        f"Checkpoint: {case.data['checkpoint']}\n"
        f"Prompt: {case.data['prompt']}\n\n"
        f"Manual review: {note}\n\n"
        f"{rollout_lines}\n\n"
        "讲解建议：先指出四路的语义差异，再比较最低分和最高分；最后强调这是组内奖励信号。",
    )


def add_section_slide(presentation: Presentation, title: str, subtitle: str, count_text: str) -> None:
    slide = blank_slide(presentation)
    add_rect(slide, 0.0, 0.0, 0.11, 7.5, TEAL, radius=False)
    add_text(slide, title, 0.78, 2.18, 11.2, 0.8, size=32, bold=True)
    add_text(slide, subtitle, 0.8, 3.08, 10.8, 0.44, size=18, color=TEAL, bold=True)
    add_rect(slide, 0.8, 3.78, 3.2, 0.05, TEAL, radius=False)
    add_text(slide, count_text, 0.8, 4.2, 10.5, 0.52, size=22, color=MUTED)
    add_text(
        slide, f"{len(presentation.slides):02d}", 12.5, 7.0, 0.42, 0.24, size=10, color=MUTED, align=PP_ALIGN.RIGHT
    )
    add_notes(slide, f"章节过渡：{title}。{subtitle}。{count_text}。")


def add_claims_slide(presentation: Presentation) -> None:
    slide = blank_slide(presentation)
    add_slide_chrome(slide, "这套材料能证明什么，也不能证明什么", "CLAIM BOUNDARY", len(presentation.slides))
    columns = (
        (
            "可以说",
            GREEN,
            [
                "完整固定合同曲线显示 RL 后总体规则分提高。",
                "同一条件的四路 rollout 具有语义差异，具备组内探索。",
                "精选案例中，奖励排序与人工观察到的完成度一致。",
                "中间 checkpoint 已能产生高低质量并存的可学习信号。",
            ],
        ),
        (
            "不要说",
            RED,
            [
                "50 个精选案例本身就是无偏的平均质量估计。",
                "不同 checkpoint 的不同样本等同于同一样本的训练时间轴。",
                "规则奖励正确就自动意味着所有语义维度都完美。",
                "单个漂亮视频可以替代 500 样本的定量评估。",
            ],
        ),
    )
    for index, (heading, color, items) in enumerate(columns):
        x = 0.55 + index * 6.3
        add_rect(slide, x, 1.35, 5.95, 4.95, PANEL, line_color=color, line_width=2.0)
        add_text(slide, heading, x + 0.3, 1.67, 5.35, 0.44, size=24, color=color, bold=True)
        lines = [(item, 16, TEXT, False) for item in items]
        add_multiline(slide, lines, x + 0.32, 2.34, 5.3, 3.3, bullet=True)
    add_text(
        slide,
        "一句话版本：曲线证明“学到了”，四路视频解释“奖励信号长什么样”。",
        0.72,
        6.68,
        11.9,
        0.35,
        size=19,
        color=TEAL,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_notes(slide, "这是建议原样保留的 claim boundary，能避免听众把高对比案例误解为随机抽样。")


def add_evidence_slide(presentation: Presentation, manifests: dict[str, Any]) -> None:
    slide = blank_slide(presentation)
    add_slide_chrome(slide, "证据链与复现入口", "AUDITED SOURCES", len(presentation.slides))
    evalkit_hash = manifests["A"]["aggregate_evidence"]["evalkit_source_sha256"]
    items = [
        ("首批 20", "80/80 exact match", "storage/eval_out/vbvr_pro_rl_rollout_demo_final_strict_20260818"),
        ("新增 30", "120/120 exact match", "storage/eval_out/vbvr_pro_rl_rollout_demo_additional30_strict_20260818"),
        ("定量曲线", "500 samples / cell", "evidence/sampler_checkpoint_trends.png + scores.csv"),
        ("构建脚本", "embedded MP4", "scripts/eval/vbvr_pro/build_rl_demo_ppt.py"),
    ]
    for index, (label, status, path) in enumerate(items):
        y = 1.33 + index * 1.22
        add_rect(slide, 0.55, y, 12.2, 0.95, PANEL, line_color=GRID)
        add_text(slide, label, 0.83, y + 0.19, 1.65, 0.3, size=17, color=TEAL, bold=True)
        add_text(slide, status, 2.62, y + 0.19, 2.35, 0.3, size=14.5, color=GREEN, bold=True)
        add_text(slide, path, 5.0, y + 0.17, 7.35, 0.34, size=11.5, color=TEXT, font=FONT_LATIN)
    add_rect(slide, 0.55, 6.35, 12.2, 0.55, PANEL_2, line_color=GRID)
    add_text(
        slide,
        f"EvalKit e140 source SHA-256: {evalkit_hash}",
        0.82,
        6.49,
        11.65,
        0.22,
        size=10.5,
        color=MUTED,
        font=FONT_LATIN,
    )
    add_notes(
        slide,
        "所有案例页的视频都直接内嵌在 PPTX 的 ppt/media 目录。manifest.json 保留 checkpoint、seed、score、"
        "native/scored SHA-256、人工复核与来源。",
    )


def build_talk_deck(
    output_path: Path,
    cases_by_id: dict[str, DemoCase],
    manifests: dict[str, Any],
    assets: dict[str, dict[str, Path]],
) -> list[DemoCase]:
    presentation = new_presentation()
    hero_cases = [cases_by_id[case_id] for case_id in HERO_IDS]
    add_cover_slide(presentation, hero_cases, assets)
    add_summary_slide(presentation, manifests)
    add_pipeline_slide(presentation)
    add_curve_slide(presentation, manifests)
    add_reading_guide_slide(presentation, cases_by_id["A02"], assets)
    for index, case in enumerate(hero_cases, start=1):
        add_case_slide(presentation, case, assets, hero_index=index)
    add_claims_slide(presentation)
    add_evidence_slide(presentation, manifests)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)
    return hero_cases


def build_full_deck(
    output_path: Path,
    cases: list[DemoCase],
    cases_by_id: dict[str, DemoCase],
    manifests: dict[str, Any],
    assets: dict[str, dict[str, Path]],
) -> list[DemoCase]:
    presentation = new_presentation()
    hero_cases = [cases_by_id[case_id] for case_id in HERO_IDS]
    remaining_cases = [case for case in cases if case.case_id not in HERO_IDS]
    add_cover_slide(presentation, hero_cases, assets)
    add_summary_slide(presentation, manifests)
    add_curve_slide(presentation, manifests)
    add_reading_guide_slide(presentation, cases_by_id["A02"], assets)
    add_section_slide(presentation, "主讲案例", "高对比、易解释、任务类型尽量多样", "10 组 · 40 个内嵌 MP4")
    for index, case in enumerate(hero_cases, start=1):
        add_case_slide(presentation, case, assets, hero_index=index)
    add_section_slide(presentation, "补充案例", "保留全部已人工复核的案例与视频", "40 组 · 160 个内嵌 MP4")
    for case in remaining_cases:
        add_case_slide(presentation, case, assets)
    add_claims_slide(presentation)
    add_evidence_slide(presentation, manifests)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)
    return hero_cases + remaining_cases


def validate_deck(path: Path, expected_slides: int, source_videos: list[Path]) -> dict[str, Any]:
    expected_hashes = Counter(sha256_file(video) for video in source_videos)
    with zipfile.ZipFile(path) as archive:
        corrupt_member = archive.testzip()
        if corrupt_member is not None:
            raise ValueError(f"Corrupt ZIP member in {path}: {corrupt_member}")
        names = archive.namelist()
        slide_names = [name for name in names if name.startswith("ppt/slides/slide") and name.endswith(".xml")]
        media_names = [name for name in names if name.startswith("ppt/media/") and name.lower().endswith(".mp4")]
        embedded_hashes = Counter(hashlib.sha256(archive.read(name)).hexdigest() for name in media_names)
        if len(slide_names) != expected_slides:
            raise ValueError(f"Expected {expected_slides} slides in {path}, found {len(slide_names)}")
        if embedded_hashes != expected_hashes:
            raise ValueError(f"Embedded MP4 hashes do not match sources in {path}")
        slide_xml = b"".join(archive.read(name) for name in slide_names)
        video_file_count = slide_xml.count(b"videoFile")
        media_element_count = slide_xml.count(b"p14:media")
        if video_file_count != len(source_videos) or media_element_count != len(source_videos):
            raise ValueError(
                f"Video OOXML mismatch in {path}: videoFile={video_file_count}, "
                f"p14:media={media_element_count}, expected={len(source_videos)}"
            )
        content_types = archive.read("[Content_Types].xml")
        if b'Extension="mp4" ContentType="video/mp4"' not in content_types:
            raise ValueError(f"Missing MP4 content type in {path}")
    reopened = Presentation(path)
    if len(reopened.slides) != expected_slides:
        raise ValueError(f"python-pptx reopen slide mismatch in {path}")
    return {
        "path": str(path.resolve()),
        "sha256": sha256_file(path),
        "bytes": path.stat().st_size,
        "slides": expected_slides,
        "embedded_mp4": len(source_videos),
        "unique_embedded_mp4_sha256": len(expected_hashes),
        "zip_test": "pass",
        "python_pptx_reopen": "pass",
        "video_ooxml": "pass",
    }


def case_video_paths(cases: list[DemoCase]) -> list[Path]:
    return [case.resolve(rollout["native_video"]) for case in cases for rollout in case.data["rollouts"]]


def write_speaker_notes(path: Path, cases: list[DemoCase], manifests: dict[str, Any]) -> None:
    curve = manifests["A"]["aggregate_evidence"]["cps_0p7"]
    lines = [
        "# VBVR-Pro RL demo speaker notes",
        "",
        "## Recommended opening",
        "",
        f"- Quantitative claim: CPS 0.7 improves from {curve['baseline_overall']:.6f} to "
        f"{curve['best_overall']:.6f} at checkpoint {curve['best_step']} "
        f"(+{curve['best_delta']:.6f}, about +{curve['best_delta'] / curve['baseline_overall']:.1%}).",
        "- Qualitative claim: each case holds the input, prompt, checkpoint, sampler, and inference contract fixed; "
        "only the rollout seed changes.",
        "- Caveat: the 50 groups are deliberately selected for contrast and are not an unbiased quality estimate.",
        "",
        "## Main ten cases",
        "",
    ]
    for index, case_id in enumerate(HERO_IDS, start=1):
        case = next(item for item in cases if item.case_id == case_id)
        scores = ", ".join(f"{score:.3f}" for score in case.data["scores"])
        seeds = ", ".join(str(item["seed"]) for item in case.data["rollouts"])
        lines.extend(
            [
                f"### {index:02d}. {case.display_label} · checkpoint {case.data['checkpoint']} · {case.task_label}",
                "",
                f"- Scores: {scores}",
                f"- Seeds: {seeds}",
                f"- Why it works: {HERO_REASONS[case.case_id]}",
                f"- Manual audit: {case.data['manual_review']['note']}",
                "- Suggested flow: play the lowest-scoring rollout, then the highest-scoring rollout, "
                "then compare the two middle cases if time allows.",
                "",
            ]
        )
    lines.extend(
        [
            "## Closing line",
            "",
            "The checkpoint curve shows that RL improves expected task reward; the four-rollout pages show the "
            "within-group behavioral contrast that makes the update possible.",
            "",
        ]
    )
    path.write_text("\n".join(lines))


def main() -> None:
    args = parse_args()
    output_dir = args.output_dir.resolve()
    output_dir.mkdir(parents=True, exist_ok=True)
    cases, manifests = load_cases()
    cases_by_id = {case.case_id: case for case in cases}
    if set(HERO_IDS) - cases_by_id.keys():
        raise ValueError("Hero case list references missing cases")
    assets = prepare_posters(cases, output_dir)

    talk_path = output_dir / "vbvr_rl_demo_talk_embedded_video_20260818.pptx"
    full_path = output_dir / "vbvr_rl_demo_full_50_embedded_video_20260818.pptx"
    talk_cases = build_talk_deck(talk_path, cases_by_id, manifests, assets)
    full_cases = build_full_deck(full_path, cases, cases_by_id, manifests, assets)

    talk_validation = validate_deck(talk_path, expected_slides=17, source_videos=case_video_paths(talk_cases))
    full_validation = validate_deck(full_path, expected_slides=58, source_videos=case_video_paths(full_cases))
    write_speaker_notes(output_dir / "speaker_notes.md", cases, manifests)

    report = {
        "schema_version": 1,
        "source_case_count": len(cases),
        "source_rollout_count": len(cases) * 4,
        "task_type_count": len({case.data["task_name"] for case in cases}),
        "checkpoints": sorted({case.data["checkpoint"] for case in cases}),
        "hero_case_ids": list(HERO_IDS),
        "talk_deck": talk_validation,
        "full_deck": full_validation,
        "evalkit_source_sha256": manifests["A"]["aggregate_evidence"]["evalkit_source_sha256"],
    }
    (output_dir / "build_report.json").write_text(json.dumps(report, ensure_ascii=False, indent=2) + "\n")
    print(json.dumps(report, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
